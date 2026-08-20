import argparse
from pathlib import Path
import pandas as pd
 
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
 
PROCESSED_DIR = Path(r"C:/Users/user/Documents/monthly-board-reporting-package/data/processed")
OUTPUT_DIR = Path(r"C:/Users/user/Documents/monthly-board-reporting-package/output")

# palette
NAVY = RGBColor(0x1F, 0x2A, 0x44)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x1B, 0x7F, 0x3B)
RED = RGBColor(0xB3, 0x26, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
 
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def load_data(month_str):
    kpi = pd.read_csv(PROCESSED_DIR / f"kpi_summary_{month_str}.csv", sep=";").iloc[0]
    variance = pd.read_csv(PROCESSED_DIR / f"variance_{month_str}.csv", sep=";")
    commentary = pd.read_csv(PROCESSED_DIR / f"commentary_{month_str}.csv", sep=";")
    forecast = pd.read_csv(PROCESSED_DIR / f"forecast_{month_str}.csv", sep=";")
    return kpi, variance, commentary, forecast
 
 
def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_title(slide, title_text, subtitle_text=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NAVY
 
    if subtitle_text:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(14)
        run2.font.color.rgb = GRAY


def add_kpi_card(slide, x, y, w, h, label, value, value_color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    shape.shadow.inherit = False
 
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
 
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(13)
    r1.font.color.rgb = GRAY
 
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = value_color
 
 
def add_footer(slide, month_str, page_note=""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    text = f"Monthly Board Reporting Package - {month_str}"
    if page_note:
        text += f"  |  {page_note}"
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY


# Title
def build_title_slide(prs, month_str):
    slide = blank_slide(prs)
 
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Monthly Board Reporting Package"
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = NAVY
 
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = f"Olist Marketplace  •  {month_str}"
    run2.font.size = Pt(20)
    run2.font.color.rgb = GRAY
 
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "Generated automatically from validated actuals, budget, and forecast data"
    run3.font.size = Pt(13)
    run3.font.color.rgb = GRAY
    run3.font.italic = True


# Financial summary
def build_summary_slide(prs, month_str, variance):
    slide = blank_slide(prs)
    add_title(slide, "Financial Summary", f"{month_str} - Actual vs. Budget")
 
    total_actual = variance["revenue"].sum()
    total_budget = variance["budget_revenue"].sum()
    total_var = total_actual - total_budget
    total_var_pct = (total_var / total_budget * 100) if total_budget else 0
    var_color = GREEN if total_var >= 0 else RED
 
    card_y = Inches(1.8)
    card_w = Inches(3.7)
    card_h = Inches(1.5)
    gap = Inches(0.25)
 
    add_kpi_card(slide, Inches(0.6), card_y, card_w, card_h, "Actual Revenue", f"${total_actual:,.0f}")
    add_kpi_card(slide, Inches(0.6) + card_w + gap, card_y, card_w, card_h, "Budget Revenue", f"${total_budget:,.0f}")
    add_kpi_card(slide, Inches(0.6) + 2 * (card_w + gap), card_y, card_w, card_h,
                 "Variance", f"${total_var:+,.0f}  ({total_var_pct:+.1f}%)", var_color)
 
    note_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(12.1), Inches(2.5))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Note: this simulated dataset's budget was built from 2016-2017 same-month-prior-year "
        "growth rates (2-10% tiers), while the underlying actuals are real 2017-2018 marketplace "
        "order data from a period of genuine hypergrowth. Large variances above largely reflect "
        "that mismatch in the data setup, not real business performance."
    )
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY
    run.font.italic = True
 
    add_footer(slide, month_str, "Financial Summary")


# KPI Dashboard
def build_kpi_slide(prs, month_str, kpi, forecast):
    slide = blank_slide(prs)
    add_title(slide, "KPI Dashboard", f"{month_str} — Totals and Trend")
 
    card_y = Inches(1.8)
    card_w = Inches(2.8)
    card_h = Inches(1.3)
    gap = Inches(0.22)
 
    rev_growth = kpi["revenue_growth_pct"]
    rev_growth_text = f"{rev_growth:+.1f}%" if pd.notna(rev_growth) else "N/A"
    rev_color = GREEN if pd.notna(rev_growth) and rev_growth >= 0 else (RED if pd.notna(rev_growth) else GRAY)
 
    add_kpi_card(slide, Inches(0.6), card_y, card_w, card_h, "Total Revenue", f"${kpi['total_revenue']:,.0f}")
    add_kpi_card(slide, Inches(0.6) + (card_w + gap), card_y, card_w, card_h, "Total Units", f"{kpi['total_units']:,.0f}")
    add_kpi_card(slide, Inches(0.6) + 2 * (card_w + gap), card_y, card_w, card_h, "Blended Avg Price", f"${kpi['blended_avg_price']:,.2f}")
    add_kpi_card(slide, Inches(0.6) + 3 * (card_w + gap), card_y, card_w, card_h, "Revenue Growth MoM", rev_growth_text, rev_color)
 
    chart_data = CategoryChartData()
    categories, actual_series, forecast_series = [], [], []
    for _, row in forecast.iterrows():
        categories.append(row["month"])
        if row["type"] == "backtest":
            actual_series.append(row["actual_revenue"])
            forecast_series.append(None)
        else:
            actual_series.append(None)
            forecast_series.append(row["forecast_revenue"])
 
    chart_data.categories = categories
    chart_data.add_series("Actual Revenue", actual_series)
    chart_data.add_series("Forecast", forecast_series)
 
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(3.5), Inches(12.1), Inches(3.3), chart_data
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue: Actual vs. Forward Forecast"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
 
    add_footer(slide, month_str, "KPI Dashboard")


# Variance Highlights
def build_variance_table_slide(prs, month_str, variance):
    slide = blank_slide(prs)
    add_title(slide, "Variance Highlights", f"{month_str} - Top Material Variances vs. Budget")
 
    material = variance[variance["material"]].copy()
    material = material.sort_values("revenue_variance", key=abs, ascending=False).head(8)
 
    rows = len(material) + 1
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.6))
    table = table_shape.table
 
    col_widths = [Inches(4.1), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
 
    headers = ["Category", "Actual", "Budget", "Variance ($)", "Variance (%)"]
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
 
    for r, (_, row) in enumerate(material.iterrows(), start=1):
        category_display = row["category"].replace("_", " ").title()
        values = [
            category_display,
            f"${row['revenue']:,.0f}",
            f"${row['budget_revenue']:,.0f}",
            f"${row['revenue_variance']:+,.0f}",
            f"{row['variance_pct']:+.1f}%" if pd.notna(row["variance_pct"]) else "N/A",
        ]
        for c, val in enumerate(values):
            cell = table.cell(r, c)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(12)
            if c == 3:
                run.font.color.rgb = GREEN if row["revenue_variance"] >= 0 else RED
 
    add_footer(slide, month_str, "Variance Highlights")


# Commentary
def build_commentary_slide(prs, month_str, commentary):
    slide = blank_slide(prs)
    add_title(slide, "Commentary", f"{month_str} - What Drove the Numbers")
 
    top_lines = commentary.sort_values("revenue_variance", key=abs, ascending=False).head(6)
 
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.1), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
 
    first = True
    for _, row in top_lines.iterrows():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = "•  " + row["narrative"]
        run.font.size = Pt(15)
        run.font.color.rgb = NAVY
        p.space_after = Pt(12)
 
    add_footer(slide, month_str, "Commentary")


# Outlook
def build_outlook_slide(prs, month_str, forecast):
    slide = blank_slide(prs)
    add_title(slide, "Outlook", f"{month_str} — Rolling Forecast")
 
    backtests = forecast[forecast["type"] == "backtest"]
    forward = forecast[forecast["type"] == "forward_forecast"].iloc[0]
 
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.1), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True
 
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "Forecast accuracy (trailing-average method, validated on known months):"
    r0.font.size = Pt(15)
    r0.font.bold = True
    r0.font.color.rgb = NAVY
 
    for _, row in backtests.iterrows():
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = (
            f"•  {row['month']}: forecast ${row['forecast_revenue']:,.0f} "
            f"vs. actual ${row['actual_revenue']:,.0f}  ({row['error_pct']:+.1f}% error)"
        )
        run.font.size = Pt(14)
        run.font.color.rgb = GRAY
        p.space_after = Pt(6)
 
    add_kpi_card(slide, Inches(0.6), Inches(3.2), Inches(4.2), Inches(1.6),
                 f"Forward Forecast — {forward['month']}", f"${forward['forecast_revenue']:,.0f}", NAVY)
 
    add_footer(slide, month_str, "Outlook")


def run(month_str):
    kpi, variance, commentary, forecast = load_data(month_str)
 
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
 
    build_title_slide(prs, month_str)
    build_summary_slide(prs, month_str, variance)
    build_kpi_slide(prs, month_str, kpi, forecast)
    build_variance_table_slide(prs, month_str, variance)
    build_commentary_slide(prs, month_str, commentary)
    build_outlook_slide(prs, month_str, forecast)
 
    out_folder = OUTPUT_DIR / month_str
    out_folder.mkdir(parents=True, exist_ok=True)
    out_path = out_folder / f"board_package_{month_str}.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
 
    return out_path
 
 
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--month", required=True, help="Target month, e.g. 2026-08")
    args = parser.parse_args()
    run(args.month)